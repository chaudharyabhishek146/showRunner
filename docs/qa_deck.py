"""Geometry and text-fit QA for the deck, without a renderer.

LibreOffice isn't installed here, so the usual render-and-look pass isn't
available. This checks analytically for the defects that pass silently and are
always user-visible: shapes off the slide, text that cannot fit its box, thin
margins, and text boxes overlapping each other.

Text fit is an estimate. Character widths are averaged per font, so treat a
small overflow as "look at this", not as proof.

    python docs/qa_deck.py docs/walkthrough-agent-deck.pptx
"""

from __future__ import annotations

import sys
from pptx import Presentation
from pptx.util import Emu

SLIDE_W = 13.333
SLIDE_H = 7.5
SAFE_MARGIN = 0.45          # flag anything closer to an edge than this
OVERFLOW_TOLERANCE = 1.06   # allow 6% before calling it overflow

# Average glyph width as a fraction of font size. Courier is monospace at
# exactly 0.6; the proportional faces are measured averages for mixed-case text.
CHAR_W = {"Courier New": 0.60, "Calibri": 0.48, "Cambria": 0.50}


def inches(v) -> float:
    return Emu(v).inches if v is not None else 0.0


def wrapped_lines(text: str, font_size: float, face: str, width: float) -> int:
    """How many lines this text needs in a box of the given width."""
    per_char = CHAR_W.get(face, 0.50) * font_size / 72.0
    if per_char <= 0 or width <= 0:
        return 1
    limit = max(1, int(width / per_char))
    lines = 0
    for hard_line in text.split("\n"):
        lines += max(1, -(-len(hard_line) // limit))  # ceil
    return lines


def check(path: str) -> int:
    prs = Presentation(path)
    problems = 0

    for index, slide in enumerate(prs.slides, start=1):
        text_boxes = []

        for shape in slide.shapes:
            x, y = inches(shape.left), inches(shape.top)
            w, h = inches(shape.width), inches(shape.height)
            label = (shape.text_frame.text[:34].replace("\n", " ")
                     if shape.has_text_frame and shape.text_frame.text
                     else shape.shape_type)

            if x < -0.01 or y < -0.01 or x + w > SLIDE_W + 0.01 or y + h > SLIDE_H + 0.01:
                print(f"  [{index}] OFF-SLIDE  {label!r} at "
                      f"({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")
                problems += 1
            elif x < SAFE_MARGIN or y < SAFE_MARGIN or \
                    x + w > SLIDE_W - SAFE_MARGIN or y + h > SLIDE_H - SAFE_MARGIN:
                print(f"  [{index}] TIGHT MARGIN  {label!r} at "
                      f"({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")
                problems += 1

            if not (shape.has_text_frame and shape.text_frame.text.strip()):
                continue

            text = shape.text_frame.text
            runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
            size = next((r.font.size.pt for r in runs if r.font.size), 14)
            face = next((r.font.name for r in runs if r.font.name), "Calibri")
            spacing = next(
                (p.line_spacing for p in shape.text_frame.paragraphs if p.line_spacing),
                None,
            )
            leading = (spacing / 12700 / 72 if isinstance(spacing, int)
                       else (spacing * size / 72 if spacing else size * 1.25 / 72))

            lines = wrapped_lines(text, size, face, w)
            needed = lines * leading
            if needed > h * OVERFLOW_TOLERANCE:
                print(f"  [{index}] OVERFLOW?  {label!r} needs ~{needed:.2f}in "
                      f"in {h:.2f}in ({lines} lines @ {size}pt {face})")
                problems += 1

            text_boxes.append((x, y, w, h, label, len(text)))

        # Text-on-text overlap. Text over a card is the design; text through
        # other text is a defect.
        for i, a in enumerate(text_boxes):
            for b in text_boxes[i + 1:]:
                ox = min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0])
                oy = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
                if ox > 0.06 and oy > 0.06:
                    print(f"  [{index}] TEXT OVERLAP  {a[4]!r} x {b[4]!r} "
                          f"({ox:.2f}in x {oy:.2f}in)")
                    problems += 1

    notes = sum(1 for s in prs.slides
                if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
    print(f"\nslides: {len(prs.slides)}   with speaker notes: {notes}")
    print(f"problems: {problems}")
    return problems


if __name__ == "__main__":
    sys.exit(1 if check(sys.argv[1]) else 0)
